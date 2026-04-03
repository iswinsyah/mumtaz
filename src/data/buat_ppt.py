from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Inisialisasi Presentasi
prs = Presentation()

def add_slide(title_text, content_points):
    slide_layout = prs.slide_layouts[1] # Layout Judul dan Isi
    slide = prs.slides.add_slide(slide_layout)
    
    # Atur Judul
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102) # Biru Tua

    # Atur Isi
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    for point in content_points:
        p = tf.add_paragraph()
        p.text = point
        p.space_after = Pt(10)
        p.level = 0

# --- SLIDE 1: JUDUL ---
slide_title_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_title_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Pilar-Pilar Pengokoh Nafsiyah Islamiyah"
subtitle.text = "Kajian: Merindukan Surga & Berlomba dalam Kebaikan\nStrategi Penguatan Jiwa Pengemban Dakwah"

# --- SLIDE 2: KONSEP NAFSIYAH ---
add_slide("Apa itu Nafsiyah Islamiyah?", [
    "Pola Sikap (Nafsiyah): Perbuatan yang didasarkan pada Akidah Islam.",
    "Bukan sekadar emosi, tapi sinkronisasi antara pikiran dan hati.",
    "Tujuan: Membentuk pribadi yang tangguh, ikhlas, dan istiqamah."
])

# --- SLIDE 3: MERINDUKAN SURGA ---
add_slide("Merindukan Surga (Syauq ilal Jannah)", [
    "Visi Melampaui Dunia: Menjadikan surga sebagai tujuan akhir setiap lelah.",
    "Penawar Ujian: Masalah dunia terasa kecil dibandingkan janji Allah.",
    "Kunci Ketangguhan: Para sahabat sanggup berkorban karena 'mencium' aroma surga sebelum syahid."
])

# --- SLIDE 4: BERLOMBA DALAM KEBAIKAN ---
add_slide("Berlomba dalam Kebaikan (Fastabiqul Khairat)", [
    "Etos Kerja Tertinggi: Tidak puas dengan amal yang 'biasa-biasa saja'.",
    "Urgensi Waktu: Menyadari ajal tidak menunggu kesiapan manusia.",
    "Prioritas Amal: Fokus pada amal yang paling dicintai Allah (Fardu)."
])

# --- SLIDE 5: PENGHALANG UTAMA ---
add_slide("Waspadai Penghalang Jiwa", [
    "Taswif (Menunda-nunda): Racun yang mematikan semangat berlomba.",
    "Wahn (Cinta Dunia): Membuat jiwa pengecut dan malas.",
    "Futur: Penurunan semangat yang harus segera diobati dengan iman."
])

# --- SLIDE 6: PENUTUP ---
add_slide("Konklusi: Menjadi Mukmin Hebat", [
    "Rindu Surga = Bahan Bakar Semangat.",
    "Berlomba dalam Kebaikan = Mesin Penggerak.",
    "Hasil: Nafsiyah yang kokoh dan tak tergoyahkan oleh fitnah zaman."
])

# Simpan File
prs.save('Presentasi_Nafsiyah_Islamiyah.pptx')